from pptx import Presentation 

prs = Presentation(path_to_presentation)

text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_textframe:
            continue
        for paragraph in shape.textframe.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)